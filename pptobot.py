import collections.abc 
from pptx import Presentation
from pptx.util import Inches



# Create presentation
prs = Presentation()

 # define slidelayouts 
slide = prs.slides.add_slide(prs.slide_layouts[0])


# title slide
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "PptoBot"
subtitle.text = " "

# Introduction page
Second_Layout = prs.slide_layouts[5]
second_slide = prs.slides.add_slide(Second_Layout)
second_slide.shapes.title.text = "Second slide"

textbox = second_slide.shapes.add_textbox(Inches(3), Inches(1.5),Inches(3), Inches(1)) 
textframe = textbox.text_frame
paragraph = textframe.add_paragraph()
paragraph.text = "This is a paragraph in the second slide!"


# end slide
slide3 = prs.slides.add_slide(prs.slide_layouts[0])
slide = slide3.shapes.title
subtitle = slide3.placeholders[1]
slide.text = "End of Presentation"
subtitle.text = "Thank You"

# export and save the presentation
prs.save('Auto_Presentation.pptx')
